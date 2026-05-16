#!/usr/bin/env python3
"""
ForsaDrive — générateur de la présentation de soutenance (version FR).

~20 diapositives, 16:9, mêmes couleurs que l'application mobile et
mêmes diagrammes que le rapport LaTeX.

Lancer :
    python3 generate_presentation_fr.py

Sortie :
    ForsaDrive_Defense_FR.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────────────────────────────────────────────────────────────
# Brand
# ─────────────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x16, 0x28)   # primary, mobile/lib/utils/app_theme.dart
NAVY_DARK   = RGBColor(0x06, 0x0E, 0x1A)
ACCENT      = RGBColor(0xE8, 0xB8, 0x4B)   # gold
ACCENT_DARK = RGBColor(0xD4, 0xA4, 0x3A)
CREAM       = RGBColor(0xFF, 0xF3, 0xCD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF3, 0xF4, 0xF6)
DARK_TEXT   = RGBColor(0x1F, 0x2A, 0x3D)
MUTED       = RGBColor(0x6B, 0x72, 0x80)

ROOT = os.path.dirname(os.path.abspath(__file__))
def img(name): return os.path.join(ROOT, name)

# 16:9 dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=DARK_TEXT,
                bullet_color=ACCENT_DARK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet marker
        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.name = 'Calibri'
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        # body
        r = p.add_run()
        r.text = line
        r.font.name = 'Calibri'
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def page_bg(slide, color=WHITE):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def add_header_bar(slide, title, eyebrow=None):
    """Top header band with section eyebrow + slide title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), NAVY)
    add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.06), ACCENT)
    if eyebrow:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(8), Inches(0.35),
                 eyebrow.upper(), size=11, bold=True, color=ACCENT)
        add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                 title, size=26, bold=True, color=WHITE)
    else:
        add_text(slide, Inches(0.5), Inches(0.32), Inches(12), Inches(0.7),
                 title, size=28, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

def add_footer(slide, page_no, total, notes=None):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), NAVY)
    add_text(slide, Inches(0.5), Inches(7.2), Inches(8), Inches(0.3),
             "ForsaDrive  •  Soutenance de fin d'études  •  ATOMIC IT",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.5), Inches(7.2), Inches(1.5), Inches(0.3),
             f"{page_no} / {total}",
             size=10, color=ACCENT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
             bold=True)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_image_centered(slide, path, top, max_w, max_h):
    if not os.path.exists(path):
        # placeholder rectangle if image missing
        ph = add_rect(slide, (SLIDE_W - max_w) / 2, top, max_w, max_h, LIGHT_GREY,
                      line=MUTED)
        add_text(slide, (SLIDE_W - max_w) / 2, top + max_h / 2 - Inches(0.2),
                 max_w, Inches(0.4), f"[image: {os.path.basename(path)}]",
                 size=12, color=MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        return ph
    pic = slide.shapes.add_picture(path, 0, top, width=max_w, height=max_h)
    # Adjust to keep aspect ratio: shrink whichever dimension exceeds the box
    if pic.width / pic.height > max_w / max_h:
        # too wide: lock width
        pic.width = max_w
        pic.height = int(max_w * (pic.height / pic.width)) if False else pic.height
    # Re-center
    pic.left = int((SLIDE_W - pic.width) / 2)
    pic.top  = top
    return pic

def add_table(slide, x, y, w, h, headers, rows, *, header_fill=NAVY,
              header_color=WHITE, body_fill=WHITE, alt_fill=LIGHT_GREY,
              text_color=DARK_TEXT, font_size=12, header_size=12):
    cols  = len(headers)
    nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, cols, x, y, w, h)
    tbl = tbl_shape.table
    # Header
    for j, h_text in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h_text
        r.font.name = 'Calibri'
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_color
    # Body
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = body_fill if i % 2 == 0 else alt_fill
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = 'Calibri'
            r.font.size = Pt(font_size)
            r.font.color.rgb = text_color
    return tbl

def add_chip(slide, x, y, label, *, fill=ACCENT, text_color=NAVY, size=11):
    w = Inches(max(0.7, 0.12 * len(label) + 0.4))
    h = Inches(0.32)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = 'Calibri'
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shp, w

# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────
TOTAL = 21  # public-facing slide count (backup slides come after this)

def slide_title():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    # decorative accent strip
    add_rect(s, 0, Inches(2.4), Inches(0.35), Inches(2.7), ACCENT)

    # Eyebrow
    add_text(s, Inches(1.0), Inches(2.3), Inches(11), Inches(0.4),
             "PROJET DE FIN D'ÉTUDES  •  2025 — 2026", size=14, bold=True,
             color=ACCENT)
    # Big title
    add_text(s, Inches(1.0), Inches(2.7), Inches(11.5), Inches(1.2),
             "ForsaDrive", size=66, bold=True, color=WHITE)
    # Subtitle
    add_text(s, Inches(1.0), Inches(3.9), Inches(11.5), Inches(0.7),
             "Plateforme de covoiturage adaptée au contexte tunisien — Web, Mobile et un backend unifié",
             size=20, color=CREAM)

    # Authors block
    add_rect(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.04), ACCENT)
    add_text(s, Inches(1.0), Inches(5.4), Inches(5.5), Inches(0.4),
             "Présenté par", size=11, color=ACCENT, bold=True)
    add_text(s, Inches(1.0), Inches(5.65), Inches(5.5), Inches(0.5),
             "Youssef BEN ABID  &  Anas YOUNES", size=18, bold=True, color=WHITE)

    add_text(s, Inches(7.0), Inches(5.4), Inches(5.5), Inches(0.4),
             "Encadré par", size=11, color=ACCENT, bold=True)
    add_text(s, Inches(7.0), Inches(5.65), Inches(5.5), Inches(0.4),
             "M. Khalil SELMI  (ATOMIC IT — encadrant professionnel)", size=13, color=WHITE)
    add_text(s, Inches(7.0), Inches(5.95), Inches(5.5), Inches(0.4),
             "Mme Ines BEN NASR  (encadrante académique)", size=13, color=WHITE)

    # Footer
    add_rect(s, 0, Inches(7.15), SLIDE_W, Inches(0.35), NAVY_DARK)
    add_text(s, Inches(0.5), Inches(7.2), Inches(12), Inches(0.3),
             "ATOMIC IT  •  Kelibia, Nabeul, Tunisia",
             size=10, color=ACCENT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

def slide_plan():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Plan de la soutenance", eyebrow="Plan")

    items = [
        ("01", "Cadre du projet & problématique", "Pourquoi la mobilité en Tunisie est un vrai problème"),
        ("02", "Solutions existantes", "Ce qui est utilisé aujourd'hui et ses limites"),
        ("03", "Solution proposée & méthodologie", "Le concept ForsaDrive et la démarche Scrum"),
        ("04", "Conception", "Architecture, cas d'utilisation, diagramme de classes"),
        ("05", "Réalisation — sprint par sprint", "Des fondations aux fonctionnalités communautaires"),
        ("06", "Tests & démo", "Ce qui fonctionne et démonstration en direct"),
        ("07", "Conclusion & perspectives", "Ce qui a été livré et la suite"),
    ]
    y = Inches(1.7)
    for num, title, sub in items:
        add_rect(s, Inches(0.6), y, Inches(0.55), Inches(0.55), ACCENT)
        add_text(s, Inches(0.6), y, Inches(0.55), Inches(0.55),
                 num, size=16, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.4), y - Inches(0.02), Inches(11.5), Inches(0.35),
                 title, size=18, bold=True, color=NAVY)
        add_text(s, Inches(1.4), y + Inches(0.28), Inches(11.5), Inches(0.3),
                 sub, size=12, color=MUTED)
        y += Inches(0.72)

    add_footer(s, 2, TOTAL)

def slide_problem():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Mobilité en Tunisie — un problème quotidien réel",
                   eyebrow="Contexte & motivation")

    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.5),
             "Le transport public est irrégulier. Les voitures privées roulent à moitié vides. "
             "Les applications existantes ne sont pas pensées pour le marché local.",
             size=15, color=DARK_TEXT)

    # Three columns of pain points
    cols = [
        ("Transport public",
         "Bus bondés, horaires irréguliers, faible couverture des villes secondaires."),
        ("Covoiturage informel",
         "Groupes Facebook, sans vérification d'identité, sans trace de paiement, sans notation."),
        ("Applications internationales",
         "BlaBlaCar/inDrive/Bolt — moyens de paiement et langue non adaptés à la Tunisie."),
    ]
    x = Inches(0.7)
    for title, body in cols:
        add_rect(s, x, Inches(2.4), Inches(4.0), Inches(2.6), LIGHT_GREY)
        add_rect(s, x, Inches(2.4), Inches(4.0), Inches(0.06), ACCENT)
        add_text(s, x + Inches(0.25), Inches(2.55), Inches(3.6), Inches(0.5),
                 title, size=16, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), Inches(3.05), Inches(3.6), Inches(1.8),
                 body, size=12, color=DARK_TEXT)
        x += Inches(4.15)

    # Quote band
    add_rect(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.4), NAVY)
    add_rect(s, Inches(0.7), Inches(5.3), Inches(0.1), Inches(1.4), ACCENT)
    add_text(s, Inches(1.0), Inches(5.45), Inches(11.5), Inches(0.4),
             "PROBLÉMATIQUE", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.9),
             "La Tunisie ne dispose pas d'une solution numérique de covoiturage "
             "sécurisée, abordable, adaptée aux habitudes de paiement locales, "
             "disponible sur le web et le mobile, et capable d'établir une vraie "
             "confiance entre conducteurs et passagers.",
             size=14, color=WHITE)

    add_footer(s, 3, TOTAL)

def slide_company():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "ATOMIC IT — entreprise d'accueil", eyebrow="Cadre du stage")

    add_bullets(s, Inches(0.7), Inches(1.7), Inches(7.5), Inches(4.5),
        [
            "Services en ingénierie informatique : développement web/mobile, conseil logiciel.",
            "Pile moderne : Linux/Windows, .NET, Java, Android/iOS, PHP.",
            "Basée à Kélibia, Nabeul. Directeur : M. Khalil SELMI.",
            "Structure plate — boucles de retour courtes entre développeurs et direction.",
            "Stage au sein de l'équipe d'ingénierie, encadré par des développeurs seniors.",
            "Scrum appliqué : stand-ups quotidiens, revues hebdomadaires avec l'encadrant.",
        ], size=15)

    # Right-side info card
    add_rect(s, Inches(8.7), Inches(1.7), Inches(4.0), Inches(4.5), LIGHT_GREY)
    add_rect(s, Inches(8.7), Inches(1.7), Inches(4.0), Inches(0.06), ACCENT)
    add_text(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(0.4),
             "EN BREF", size=11, bold=True, color=NAVY)

    info = [
        ("Activité", "Développement web & mobile"),
        ("Directeur", "M. Khalil SELMI"),
        ("Lieu", "Kélibia, Nabeul"),
        ("Téléphone", "+216 55 343 224"),
        ("Email", "contact@atomicitpro.com"),
        ("Méthodologie", "Agile / Scrum"),
    ]
    y = Inches(2.35)
    for k, v in info:
        add_text(s, Inches(8.95), y, Inches(1.3), Inches(0.32),
                 k, size=10, bold=True, color=MUTED)
        add_text(s, Inches(10.25), y, Inches(2.4), Inches(0.32),
                 v, size=11, color=DARK_TEXT)
        y += Inches(0.36)

    add_footer(s, 4, TOTAL)

def slide_existing():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Ce qui est utilisé aujourd'hui", eyebrow="Étude critique")

    headers = ["Critère", "Facebook", "inDrive", "Bolt", "BlaBlaCar", "ForsaDrive"]
    rows = [
        ["Vérification d'identité",        "—",   "Partielle", "✓", "✓", "✓"],
        ["Paiement intégré",               "—",   "—",         "✓", "✓", "✓"],
        ["Trajets longue distance",        "✓",   "—",         "—", "✓", "✓"],
        ["Système de notation",            "—",   "✓",         "✓", "✓", "✓"],
        ["Adapté au contexte local",       "✓",   "Partiel",   "Partiel", "—", "✓"],
        ["Réduction étudiante (50%)",      "—",   "—",         "—", "—", "✓"],
        ["Portefeuille interne",           "—",   "—",         "—", "—", "✓"],
    ]
    add_table(s, Inches(0.7), Inches(1.7), Inches(12), Inches(4.6),
              headers, rows, font_size=12, header_size=12)

    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
             "Aucune solution existante ne combine l'adaptation locale et des mécanismes "
             "de confiance structurés. ForsaDrive se positionne précisément dans ce vide.",
             size=13, bold=True, color=NAVY)

    add_footer(s, 5, TOTAL)

def slide_solution():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "ForsaDrive — ce que nous avons construit", eyebrow="Solution proposée")

    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.6),
        "Un écosystème intégré : application web, application mobile et un backend partagé, "
        "couvrant tout le cycle de vie, de l'inscription jusqu'à la notation post-trajet.",
        size=14, color=DARK_TEXT)

    pillars = [
        ("Adaptation locale",
         "Acompte de 50 % + règlement en cash, français/anglais/arabe avec RTL, "
         "domaines universitaires tunisiens reconnus."),
        ("Confiance",
         "Vérification d'identité et de statut étudiant, notation des deux côtés, "
         "gestion des plaintes, score de fiabilité."),
        ("Fonctions intelligentes",
         "Score de compatibilité, tableau de bord conducteur, mise en avant des trajets, "
         "fil social, messagerie intégrée, helpdesk assisté."),
    ]
    x = Inches(0.7)
    for title, body in pillars:
        add_rect(s, x, Inches(2.5), Inches(4.0), Inches(3.0), NAVY)
        add_rect(s, x, Inches(2.5), Inches(4.0), Inches(0.08), ACCENT)
        add_text(s, x + Inches(0.25), Inches(2.65), Inches(3.6), Inches(0.5),
                 title, size=18, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.25), Inches(3.15), Inches(3.6), Inches(2.2),
                 body, size=13, color=WHITE)
        x += Inches(4.15)

    # KPI row
    kpis = [
        ("3", "applications coordonnées"),
        ("4", "sprints fonctionnels"),
        ("20+", "entités métier"),
        ("3", "langues supportées"),
    ]
    x = Inches(0.7)
    for v, lbl in kpis:
        add_rect(s, x, Inches(5.7), Inches(2.95), Inches(1.05), CREAM)
        add_text(s, x, Inches(5.75), Inches(2.95), Inches(0.55),
                 v, size=28, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(6.25), Inches(2.95), Inches(0.4),
                 lbl, size=11, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER)
        x += Inches(3.05)

    add_footer(s, 6, TOTAL)

def slide_methodology():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Scrum — cinq itérations en deux releases",
                   eyebrow="Méthodologie")

    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.5),
             "Le Sprint 0 a préparé le terrain. Quatre sprints fonctionnels ont ensuite "
             "livré deux releases cohérentes du produit.",
             size=14, color=DARK_TEXT)

    # Sprint timeline (5 boxes)
    sprints = [
        ("S0", "Architecture\n& Conception", NAVY, ACCENT),
        ("S1", "Fondations\n(Auth + Profils)", ACCENT, NAVY),
        ("S2", "Trajets\n& Réservations", ACCENT, NAVY),
        ("S3", "Paiements\n& Intelligence", ACCENT, NAVY),
        ("S4", "Communauté\n& Finalisation", ACCENT, NAVY),
    ]
    box_w = Inches(2.3)
    gap   = Inches(0.15)
    total_w = box_w * 5 + gap * 4
    x = (SLIDE_W - total_w) / 2
    y = Inches(2.6)
    for code, title, fill, text_color in sprints:
        add_rect(s, x, y, box_w, Inches(2.0), fill)
        add_text(s, x, y + Inches(0.15), box_w, Inches(0.5),
                 code, size=22, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.85), box_w - Inches(0.2),
                 Inches(1.0), title, size=12, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += box_w + gap

    # Release labels
    add_rect(s, Inches(3.0), Inches(4.85), Inches(4.6), Inches(0.04), ACCENT)
    add_text(s, Inches(3.0), Inches(4.95), Inches(4.6), Inches(0.4),
             "RELEASE 1 — boucle de réservation utilisable",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(7.7), Inches(4.85), Inches(4.6), Inches(0.04), ACCENT)
    add_text(s, Inches(7.7), Inches(4.95), Inches(4.6), Inches(0.4),
             "RELEASE 2 — intelligence + communauté",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Process bullets
    add_bullets(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.4),
        [
            "Stand-ups quotidiens au sein de l'équipe — revues hebdomadaires avec l'encadrant.",
            "Chaque sprint a produit un incrément utilisable, revu et démontré.",
            "Tableau Trello pour garder le backlog visible — story points sur une échelle Fibonacci.",
        ], size=13)

    add_footer(s, 7, TOTAL)

def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Architecture — client / serveur, trois couches",
                   eyebrow="Sprint 0 · Conception")

    # Left column — bullets
    add_bullets(s, Inches(0.7), Inches(1.7), Inches(5.4), Inches(5.0),
        [
            "Deux clients (web + mobile) consomment une seule API REST.",
            "Backend en PHP, exposé sous /api/, sécurisé par bearer tokens.",
            "Trois couches : présentation, application, données.",
            "Services externes (paiement, notifications, stockage) joignables "
            "uniquement via le backend — aucun secret côté client.",
            "HTTPS en production. Mots de passe bcrypt. Journal d'audit sur les "
            "opérations sensibles.",
        ], size=14)

    # Right side — schematic boxes
    box_x = Inches(7.0); box_y = Inches(1.85)
    add_rect(s, box_x,             box_y,             Inches(2.6), Inches(0.9), CREAM)
    add_text(s, box_x, box_y, Inches(2.6), Inches(0.9), "Application Web",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, box_x + Inches(3.0), box_y,             Inches(2.6), Inches(0.9), CREAM)
    add_text(s, box_x + Inches(3.0), box_y, Inches(2.6), Inches(0.9), "Application Mobile",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, box_x + Inches(1.5), box_y + Inches(1.4), Inches(2.6), Inches(0.95), NAVY)
    add_text(s, box_x + Inches(1.5), box_y + Inches(1.4), Inches(2.6), Inches(0.95),
             "Backend\n(API REST)",
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, box_x + Inches(1.5), box_y + Inches(2.85), Inches(2.6), Inches(0.85), ACCENT)
    add_text(s, box_x + Inches(1.5), box_y + Inches(2.85), Inches(2.6), Inches(0.85),
             "SQLite (WAL)",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Connecting lines (visual hint)
    line1 = s.shapes.add_connector(2, box_x + Inches(1.3), box_y + Inches(0.9),
                                   box_x + Inches(2.5), box_y + Inches(1.4))
    line1.line.color.rgb = NAVY
    line2 = s.shapes.add_connector(2, box_x + Inches(4.3), box_y + Inches(0.9),
                                   box_x + Inches(3.0), box_y + Inches(1.4))
    line2.line.color.rgb = NAVY
    line3 = s.shapes.add_connector(2, box_x + Inches(2.8), box_y + Inches(2.35),
                                   box_x + Inches(2.8), box_y + Inches(2.85))
    line3.line.color.rgb = NAVY

    add_footer(s, 8, TOTAL)

def slide_class_diagram():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Diagramme de classes — 20+ entités, 5 groupes",
                   eyebrow="Sprint 0 · Conception")

    add_image_centered(s, img("forsadrive_class_diagram.png"),
                       Inches(1.4), Inches(11.5), Inches(4.95))

    # Caption / groups (above the footer band at 7.15)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             "Utilisateurs & profils  ·  Trajets & réservations  ·  Paiements & notations  "
             "·  Vérification & réductions  ·  Communication & modération",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(s, 9, TOTAL)

def slide_use_case():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Cas d'utilisation global — quatre acteurs", eyebrow="Sprint 0 · Conception")
    add_image_centered(s, img("ForsaDrive_UseCase.png"),
                       Inches(1.4), Inches(11.5), Inches(4.95))
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             "Passager  ·  Conducteur  ·  Administrateur  ·  Système (automatique)",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, 10, TOTAL)

def slide_sprint_section(num, color_dark, color_accent, title, mission,
                         backlog_rows, deliverables, page_no):
    """Section divider + content side by side, single slide per sprint."""
    s = prs.slides.add_slide(BLANK)
    page_bg(s)

    # Left rail (sprint marker)
    add_rect(s, 0, 0, Inches(3.4), SLIDE_H, color_dark)
    add_text(s, Inches(0.4), Inches(0.6), Inches(2.6), Inches(0.5),
             f"SPRINT {num}", size=14, bold=True, color=color_accent)
    # MISSION label is rendered below from the sprint definition
    add_text(s, Inches(0.4), Inches(1.0), Inches(2.6), Inches(2.0),
             title, size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.4), Inches(3.4), Inches(2.6), Inches(0.4),
             "OBJECTIF", size=11, bold=True, color=color_accent)
    add_text(s, Inches(0.4), Inches(3.7), Inches(2.7), Inches(2.0),
             mission, size=12, color=WHITE)

    # Right area — backlog table
    add_text(s, Inches(3.7), Inches(0.4), Inches(9), Inches(0.5),
             "Backlog du sprint (extrait)", size=14, bold=True, color=NAVY)
    headers = ["#", "User story", "Pts"]
    add_table(s, Inches(3.7), Inches(0.85), Inches(9.3), Inches(3.3),
              headers, backlog_rows, font_size=11, header_size=11)

    # Deliverables row
    add_text(s, Inches(3.7), Inches(4.35), Inches(9), Inches(0.4),
             "Livrables", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(3.7), Inches(4.75), Inches(9.3), Inches(2.3),
                deliverables, size=12)

    add_footer(s, page_no, TOTAL)

def slide_sprint1():
    slide_sprint_section(
        1, NAVY, ACCENT,
        "Fondations",
        "Création de compte, authentification, profil, vérification "
        "étudiante, candidature conducteur et validation par l'admin.",
        [
            ["US1.1", "Créer un compte",                            "5"],
            ["US1.2", "Se connecter avec email + mot de passe",     "3"],
            ["US1.3", "Voir / modifier le profil",                  "3"],
            ["US1.4", "Vérification étudiante par OTP universitaire","8"],
            ["US1.5", "Postuler comme conducteur (permis + véhicule)","8"],
            ["US1.6", "L'admin examine les candidatures conducteur","5"],
            ["US1.7", "L'admin gère les domaines universitaires",   "3"],
        ],
        [
            "Flux OTP en libre-service (code 10 min, 5 tentatives max, 30 s entre envois).",
            "À l'approbation, création automatique de DriverProfile + Véhicule.",
            "Entrée dans le journal d'audit à chaque vérification.",
        ],
        11)

def slide_sprint2():
    slide_sprint_section(
        2, NAVY, ACCENT,
        "Trajets & Réservations",
        "Un conducteur publie un trajet. Un passager cherche, filtre, "
        "réserve avec un acompte de 50 %.",
        [
            ["US2.1", "Publier un trajet (origine, destination, date, prix, places)", "5"],
            ["US2.2", "Modifier ou annuler un trajet",                                "3"],
            ["US2.3", "Rechercher des trajets par origine / destination / date",      "5"],
            ["US2.4", "Filtrer par prix, heure de départ, note du conducteur",        "5"],
            ["US2.5", "Voir le détail du trajet et le profil du conducteur",          "3"],
            ["US2.6", "Réserver des places avec acompte de 50 %",                     "8"],
            ["US2.7", "Réservation de groupe pour des amis",                          "5"],
            ["US2.8", "Le conducteur accepte / refuse les demandes",                  "5"],
        ],
        [
            "La réservation n'est créée qu'après succès de l'acompte — pas de réservation orpheline.",
            "Index sur from_location, to_location et departure_time pour la recherche.",
            "Le diagramme d'activité couvre la collecte cash en fin de trajet.",
        ],
        12)

def slide_sprint3():
    slide_sprint_section(
        3, NAVY, ACCENT,
        "Paiements & Intelligence",
        "Portefeuille, codes promo, mise en avant de trajets, "
        "tableau de bord conducteur, score de compatibilité.",
        [
            ["US3.1", "Voir solde et historique du portefeuille",          "5"],
            ["US3.2", "Recharger le portefeuille",                          "3"],
            ["US3.3", "Appliquer auto. la réduction étudiante (50 %)",     "3"],
            ["US3.4", "Appliquer un code promo organisationnel",            "3"],
            ["US3.5", "Mettre un trajet en avant (boost)",                  "5"],
            ["US3.6", "Tableau de bord : fiabilité + revenus + notes",     "8"],
            ["US3.7", "Badge de score de compatibilité par trajet",        "5"],
            ["US3.8", "Recalculer la fiabilité après chaque trajet terminé","3"],
        ],
        [
            "Pipeline prix : base × places → étudiant 50 % → promo % → acompte 50 %.",
            "Score de fiabilité : note × taux de complétion × (1 − taux d'annulation).",
            "Bug détecté en revue : la promo s'appliquait avant la réduction étudiante — corrigé.",
        ],
        13)

def slide_sprint4():
    slide_sprint_section(
        4, NAVY, ACCENT,
        "Communauté & Finalisation",
        "Messagerie, notations, plaintes, fil social, helpdesk, "
        "interface multilingue, panneau d'administration.",
        [
            ["US4.1", "Messagerie en temps réel (polling)",          "5"],
            ["US4.2", "Notifications push (FCM)",                     "3"],
            ["US4.3", "Fil social (posts, likes, commentaires)",      "8"],
            ["US4.4", "Noter l'autre partie après un trajet",         "3"],
            ["US4.5", "Déposer une plainte après un trajet",          "5"],
            ["US4.6", "HelpDesk assisté avec escalade vers humain",   "5"],
            ["US4.7", "Changer la langue FR / EN / AR (RTL)",         "5"],
            ["US4.8", "Panneau admin — 7 onglets",                    "8"],
        ],
        [
            "Polling toutes les 3 s — suffisant à notre échelle, sans dépendance WebSocket.",
            "Bot HelpDesk : 14 catégories FAQ, escalade vers un agent humain si besoin.",
            "L'arabe bascule automatiquement la mise en page via les localizations Flutter.",
        ],
        14)

def slide_screens():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Quelques écrans", eyebrow="Réalisation")

    # 4-column grid of mobile screens
    captions = [
        ("Authentification", "mobile_auth.png"),
        ("Recherche + score", "mobile_search.png"),
        ("Détail trajet + carte", "mobile_trip_details.png"),
        ("Tableau de bord conducteur", "mobile_driver_dashboard.png"),
    ]
    x = Inches(0.5); y = Inches(1.6)
    cell_w = Inches(3.0)
    for caption, fname in captions:
        add_rect(s, x, y, cell_w, Inches(4.4), LIGHT_GREY)
        add_image_centered_into(s, img(fname), x + Inches(0.1), y + Inches(0.1),
                                cell_w - Inches(0.2), Inches(4.0))
        add_text(s, x, y + Inches(4.1), cell_w, Inches(0.3),
                 caption, size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        x += cell_w + Inches(0.13)

    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             "Les équivalents web couvrent les mêmes parcours pour les utilisateurs "
             "desktop et le panneau d'administration.",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(s, 15, TOTAL)

def add_image_centered_into(slide, path, x, y, w, h):
    """Like add_image_centered but inside an arbitrary cell box."""
    if not os.path.exists(path):
        ph = add_rect(slide, x, y, w, h, WHITE, line=MUTED)
        add_text(slide, x, y + h / 2 - Inches(0.2), w, Inches(0.4),
                 f"[{os.path.basename(path)}]", size=10, color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return ph
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    # Center inside cell
    pic.left = int(x + (w - pic.width) / 2)
    pic.top  = int(y + (h - pic.height) / 2)
    return pic

def slide_mobile_features():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Ce que seul le mobile permet",
                   eyebrow="Fonctionnalités mobiles")

    cells = [
        ("Notifications push",
         "Firebase Cloud Messaging livre les alertes de réservation, message et trajet "
         "même quand l'app est en arrière-plan."),
        ("Carte interactive",
         "Tuiles OpenStreetMap + routage OSRM pour tracer le trajet entre "
         "origine et destination sur chaque écran de détail."),
        ("Messagerie temps réel",
         "Conversations liées aux réservations — bulles, séparateurs par date, accusés de lecture. "
         "Polling 3 s, sans dépendance WebSocket."),
        ("Multilingue + RTL",
         "FR / EN / AR avec bascule automatique en RTL pour l'arabe, "
         "sans aucun override manuel dans les widgets."),
    ]
    x = Inches(0.6); y = Inches(1.6)
    w = Inches(6.0); h = Inches(2.55)
    for i, (title, body) in enumerate(cells):
        col = i % 2; row = i // 2
        cx = Inches(0.6) + col * (Inches(6.15))
        cy = Inches(1.6) + row * (Inches(2.7))
        add_rect(s, cx, cy, w, h, LIGHT_GREY)
        add_rect(s, cx, cy, Inches(0.08), h, ACCENT)
        add_text(s, cx + Inches(0.25), cy + Inches(0.2), w - Inches(0.4),
                 Inches(0.5), title, size=18, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.25), cy + Inches(0.85), w - Inches(0.4),
                 Inches(1.6), body, size=12, color=DARK_TEXT)

    add_footer(s, 16, TOTAL)

def slide_stack():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Pile technologique", eyebrow="Réalisation")

    # Three columns — frontend / backend / mobile
    blocks = [
        ("Frontend web", NAVY, [
            ("PHP 8 (rendu serveur)", "Routage, sessions, rendu des pages."),
            ("HTML5 + Bootstrap 5", "Mise en page responsive, composants cohérents."),
            ("JavaScript vanille", "Validation formulaires, AJAX, onglets admin."),
        ]),
        ("Backend / données", ACCENT_DARK, [
            ("API REST PHP", "Front-controller sous /api/, réponses JSON."),
            ("PDO + bearer tokens", "Requêtes paramétrées, tokens hex 32 octets."),
            ("SQLite (mode WAL)", "Fichier partagé entre web et mobile."),
        ]),
        ("Mobile", NAVY, [
            ("Flutter (Dart)", "Un seul code, rendu natif sur Android et iOS."),
            ("Provider + go_router", "Gestion d'état + navigation déclarative."),
            ("FCM, flutter_map, l10n", "Push, carte OSM, FR/EN/AR avec RTL."),
        ]),
    ]

    x = Inches(0.6); y = Inches(1.6)
    w = Inches(4.0); h = Inches(5.0)
    for title, head_color, items in blocks:
        add_rect(s, x, y, w, Inches(0.65), head_color)
        add_text(s, x, y, w, Inches(0.65), title,
                 size=16, bold=True,
                 color=WHITE if head_color != ACCENT_DARK else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, y + Inches(0.65), w, h - Inches(0.65), LIGHT_GREY)
        cy = y + Inches(0.85)
        for tech, sub in items:
            add_text(s, x + Inches(0.25), cy, w - Inches(0.4), Inches(0.35),
                     tech, size=13, bold=True, color=NAVY)
            add_text(s, x + Inches(0.25), cy + Inches(0.35),
                     w - Inches(0.4), Inches(0.7),
                     sub, size=11, color=DARK_TEXT)
            cy += Inches(1.25)
        x += Inches(4.2)

    add_footer(s, 17, TOTAL)

def slide_tests():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Tests — en continu, pas reportés",
                   eyebrow="Qualité")

    # Left: testing strategy
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.4),
             "Stratégie", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.0),
        [
            "Tests unitaires sur le pipeline de prix, l'acompte, la machine à états des réservations.",
            "Collections Postman rejouant les scénarios bout-en-bout après chaque changement.",
            "Tests manuels multi-appareils sur trois téléphones Android.",
            "Démo hebdomadaire avec l'encadrant pour repérer les soucis d'utilisabilité.",
        ], size=12)

    # Right: results table
    add_text(s, Inches(6.5), Inches(1.7), Inches(6.3), Inches(0.4),
             "Résultats", size=16, bold=True, color=NAVY)
    rows = [
        ["Inscription + OTP étudiant",         "✓ Réussi"],
        ["Session par bearer token",            "✓ Réussi"],
        ["Candidature conducteur + approbation","✓ Réussi"],
        ["Publication / annulation d'un trajet","✓ Réussi"],
        ["Réservation avec acompte 50 %",       "✓ Réussi"],
        ["Réduction étudiante + code promo",    "✓ Réussi"],
        ["Recharge + historique portefeuille",  "✓ Réussi"],
        ["Notations + plaintes",                "✓ Réussi"],
        ["Chat temps réel + push (FCM)",        "✓ Réussi"],
        ["Multilingue + RTL",                   "✓ Réussi"],
    ]
    add_table(s, Inches(6.5), Inches(2.1), Inches(6.3), Inches(4.7),
              ["Scénario", "Résultat"], rows, font_size=11, header_size=11)

    add_footer(s, 18, TOTAL)

def slide_demo():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    add_rect(s, 0, Inches(3.4), Inches(0.4), Inches(0.8), ACCENT)
    add_text(s, Inches(1.0), Inches(2.5), Inches(11.5), Inches(0.5),
             "DÉMONSTRATION", size=18, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(3.0), Inches(11.5), Inches(1.5),
             "ForsaDrive en action", size=64, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.7), Inches(11.5), Inches(0.7),
             "Parcours : inscription → vérification étudiante → recherche → "
             "réservation → acceptation conducteur → collecte cash.",
             size=18, color=CREAM)
    add_footer(s, 19, TOTAL)

def slide_limitations():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Limites & feuille de route", eyebrow="Évaluation honnête")

    # Two columns
    add_rect(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.0), LIGHT_GREY)
    add_rect(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(0.06), ACCENT)
    add_text(s, Inches(0.85), Inches(1.75), Inches(5.5), Inches(0.5),
             "Ce qui manque aujourd'hui", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(5.6), Inches(4.5),
        [
            "Le portefeuille fonctionne par rechargement manuel — pas encore de vraie passerelle de paiement.",
            "SQLite convient au développement, pas à une production à fort trafic.",
            "Mobile testé uniquement sur Android — build iOS non vérifié.",
            "Pas encore de flux de mot de passe oublié par email.",
            "Les performances sous charge concurrente n'ont pas été mesurées.",
        ], size=12, bullet_color=RGBColor(0xC0, 0x39, 0x2B))

    add_rect(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(5.0), CREAM)
    add_rect(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(0.06), NAVY)
    add_text(s, Inches(6.95), Inches(1.75), Inches(5.5), Inches(0.5),
             "Étapes suivantes", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(6.95), Inches(2.3), Inches(5.6), Inches(4.5),
        [
            "Migrer vers PostgreSQL ou MySQL pour la production.",
            "Intégrer une passerelle de paiement tunisienne (Konnect ou Paymee).",
            "Déployer le backend en HTTPS sur un VPS, sauvegardes automatisées.",
            "Publier l'application Flutter sur le Google Play Store.",
            "Ajouter 2FA pour conducteurs/admins + cache hors-ligne sur mobile.",
        ], size=12)

    add_footer(s, 20, TOTAL)

def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    add_rect(s, 0, Inches(3.0), Inches(0.4), Inches(1.5), ACCENT)
    add_text(s, Inches(1.0), Inches(2.6), Inches(11.5), Inches(0.5),
             "QUESTIONS & DISCUSSION", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(3.1), Inches(11.5), Inches(1.5),
             "Merci pour votre attention.",
             size=58, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.7), Inches(11.5), Inches(0.6),
             "Nous sommes prêts à répondre à vos questions.",
             size=20, color=CREAM)

    # Bottom names band
    add_rect(s, 0, Inches(6.5), SLIDE_W, Inches(0.65), NAVY_DARK)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.55),
             "Youssef BEN ABID  &  Anas YOUNES   ·   ForsaDrive   ·   2025 — 2026",
             size=14, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             bold=True)
    add_footer(s, 21, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────────────────────────────────────
slide_title()
slide_plan()
slide_problem()
slide_company()
slide_existing()
slide_solution()
slide_methodology()
slide_architecture()
slide_class_diagram()
slide_use_case()
slide_sprint1()
slide_sprint2()
slide_sprint3()
slide_sprint4()
slide_screens()
slide_mobile_features()
slide_stack()
slide_tests()
slide_demo()
slide_limitations()
slide_thanks()

# ─────────────────────────────────────────────────────────────────────────────
# Backup slides — keep at the end of the deck, navigate to them on demand
# during the Q&A if the jury asks for deeper detail.
# ─────────────────────────────────────────────────────────────────────────────
def slide_backup_divider():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    add_rect(s, 0, Inches(3.0), Inches(0.35), Inches(1.5), ACCENT)
    add_text(s, Inches(1.0), Inches(2.6), Inches(11), Inches(0.4),
             "DIAPOSITIVES DE SECOURS", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(3.1), Inches(11.5), Inches(1.2),
             "Annexes — pour les questions",
             size=44, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.7), Inches(11.5), Inches(0.5),
             "Les diapositives suivantes ne font pas partie de la présentation principale. "
             "Elles sont là au cas où le jury demanderait plus de détails.",
             size=14, color=CREAM)

def slide_backup_seq_book():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Séquence — réservation par un passager", eyebrow="Annexe A · Flux de réservation")
    add_image_centered(s, img("ForsaDrive_Sequence_BookRide.png"),
                       Inches(1.4), Inches(11.5), Inches(5.4))
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
             "L'acompte est débité du portefeuille AVANT que la réservation ne soit créée "
             "— garantit qu'il n'y a pas de réservation orpheline.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, "A", TOTAL)

def slide_backup_seq_otp():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Séquence — vérification étudiante par OTP", eyebrow="Annexe B · Confiance")
    add_image_centered(s, img("ForsaDrive_Sequence_VerifyStudent.png"),
                       Inches(1.4), Inches(11.5), Inches(5.4))
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
             "Domaine vérifié dans student_domains • OTP à 6 chiffres, 10 min d'expiration, "
             "5 tentatives max, 30 s entre deux envois.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, "B", TOTAL)

def slide_backup_seq_driver():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Séquence — candidature conducteur & approbation admin",
                   eyebrow="Annexe C · Onboarding")
    add_image_centered(s, img("ForsaDrive_Sequence_DriverApplication.png"),
                       Inches(1.4), Inches(11.5), Inches(5.4))
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
             "À l'approbation : DriverProfile + Véhicule sont créés et le rôle utilisateur "
             "passe à DRIVER de manière atomique.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, "C", TOTAL)

def slide_backup_activity():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Diagramme d'activité — parcours complet de réservation",
                   eyebrow="Annexe D · Cycle complet")
    add_image_centered(s, img("ForsaDrive_Activity_BookingFlow.png"),
                       Inches(1.4), Inches(11.5), Inches(5.4))
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
             "Trois couloirs : passager / système / conducteur. Décisions : réessayer "
             "le paiement, accepter/refuser, départ, arrivée, collecte cash.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, "D", TOTAL)

def slide_backup_price_pipeline():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Pipeline de calcul du prix", eyebrow="Annexe E · Règle métier")

    add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
             "L'ordre compte. Revu et validé avec l'encadrant.",
             size=14, color=DARK_TEXT)

    steps = [
        ("1", "prix de base",       "price_per_seat × seats_booked"),
        ("2", "réduction étudiante","si is_student_verified : price × 0.5"),
        ("3", "code promo",         "si code org valide : price × (1 − pct/100)"),
        ("4", "acompte",            "prepayment = price × 0.5"),
        ("5", "persister",          "stocker les 5 valeurs sur la ligne booking"),
    ]
    y = Inches(2.4)
    for num, label, formula in steps:
        add_rect(s, Inches(1.0), y, Inches(0.5), Inches(0.5), ACCENT)
        add_text(s, Inches(1.0), y, Inches(0.5), Inches(0.5), num,
                 size=16, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.7), y + Inches(0.02), Inches(3), Inches(0.35),
                 label, size=14, bold=True, color=NAVY)
        add_text(s, Inches(5.0), y + Inches(0.05), Inches(7.5), Inches(0.4),
                 formula, size=13, color=DARK_TEXT, font='Consolas')
        y += Inches(0.7)

    add_rect(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.5), CREAM)
    add_text(s, Inches(0.85), Inches(6.35), Inches(11.8), Inches(0.4),
             "Bug détecté en revue de Sprint 3 : étapes 2 et 3 inversées. "
             "Corrigé le soir même, verrouillé par un test unitaire.",
             size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, "E", TOTAL)

def slide_backup_db_schema():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Schéma relationnel principal (extrait)", eyebrow="Annexe F · Modèle de données")

    rows = [
        ["users",          "id, username, email, password, is_driver, is_student, balance, score"],
        ["driver_profiles","id, user_id*, license_number, avg_rating, total_trips, reliability"],
        ["vehicles",       "id, user_id*, type, make, model, plate_number, seats, has_ac, has_wifi"],
        ["rides",          "id, driver_id*, vehicle_id*, from_location, to_location, departure_time, price, available_seats, status"],
        ["bookings",       "id, ride_id*, passenger_id*, seats, paid_amount, status"],
        ["payments",       "id, user_id*, amount, type, description, ref_id"],
        ["ratings",        "id, ride_id*, from_user_id*, to_user_id*, score, comment"],
        ["student_verifications", "id, user_id*, university_email, status, verified_at"],
        ["organizations",  "id, name, type, contact_email, email_domain, discount_percent, discount_code, status"],
        ["complaints",     "id, from_user_id*, against_user_id*, type, text, status, admin_note"],
        ["messages",       "id, conversation_id*, sender_id*, body, is_read, sent_at"],
    ]
    add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.8),
              ["Table", "Colonnes"], rows, font_size=10, header_size=11)

    add_text(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.4),
             "* clé étrangère. ~20 tables au total dans le schéma. "
             "Index sur origin, destination, departure_time et clés étrangères.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, "F", TOTAL)

def slide_backup_api():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Surface de l'API REST (extrait)", eyebrow="Annexe G · Contrat API")

    rows = [
        ["POST",   "/api/auth/register",            "Créer un compte"],
        ["POST",   "/api/auth/login",               "Retourne un bearer token de 32 octets"],
        ["GET",    "/api/rides",                    "Rechercher avec from/to/date/filtres"],
        ["POST",   "/api/rides",                    "Le conducteur publie un trajet"],
        ["DELETE", "/api/rides/{id}",               "Le conducteur annule"],
        ["POST",   "/api/bookings",                 "Le passager réserve — acompte 50 %"],
        ["POST",   "/api/bookings/group",           "Réservation de groupe"],
        ["POST",   "/api/bookings/validate-promo",  "Vérification d'un code promo organisationnel"],
        ["POST",   "/api/bookings/requests/{id}",   "Le conducteur accepte/refuse"],
        ["POST",   "/api/student/send-otp",         "Envoyer un OTP à l'email universitaire"],
        ["POST",   "/api/student/verify-otp",       "Vérifier l'OTP, marquer vérifié"],
        ["GET",    "/api/admin/organizations",      "Lister les organisations (admin)"],
        ["POST",   "/api/admin/organizations/{id}/review", "Approuver / rejeter une organisation"],
        ["GET",    "/api/feed",                     "Fil social (posts + interactions)"],
    ]
    add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.0),
              ["Méthode", "Chemin", "Rôle"], rows, font_size=11, header_size=11)

    add_text(s, Inches(0.5), Inches(6.75), Inches(12), Inches(0.3),
             "Tous les endpoints protégés requièrent l'en-tête Authorization: Bearer <token>.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, "G", TOTAL)

def slide_backup_security():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Posture de sécurité", eyebrow="Annexe H · Risques & mitigations")

    rows = [
        ["Stockage mots de passe",  "bcrypt via PHP password_hash"],
        ["Session",                 "Bearer token aléatoire 32 octets hex, TTL 30 jours"],
        ["Transport",               "HTTPS en production (HTTP via XAMPP en dev)"],
        ["Injection SQL",           "PDO avec requêtes paramétrées partout"],
        ["XSS",                     "Sortie serveur échappée (htmlspecialchars)"],
        ["CSRF (admin web)",        "Panneau admin same-origin + session"],
        ["Séparation admin",        "Connexion admin bloquée sur l'API mobile (web uniquement)"],
        ["Téléversement fichiers",  "Validation mime + taille ; servi depuis /Src sous XAMPP"],
        ["Journal d'audit",         "Opérations sensibles (vérif., suspension) tracées côté serveur"],
        ["Limitation de débit",     "Envoi OTP : cooldown de 30 s par compte"],
    ]
    add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3),
              ["Préoccupation", "Mitigation"], rows, font_size=12, header_size=12)

    add_footer(s, "H", TOTAL)

# Build backup deck
slide_backup_divider()
slide_backup_seq_book()
slide_backup_seq_otp()
slide_backup_seq_driver()
slide_backup_activity()
slide_backup_price_pipeline()
slide_backup_db_schema()
slide_backup_api()
slide_backup_security()

# ─────────────────────────────────────────────────────────────────────────────
# Speaker notes (visible in PowerPoint presenter view, hidden from audience)
# ─────────────────────────────────────────────────────────────────────────────
SPEAKER_NOTES = [
    # 1 — Title
    "Bonjour, je suis Youssef, voici Anas. Nous allons vous présenter "
    "ForsaDrive, notre PFE réalisé chez ATOMIC IT sous la supervision de "
    "M. Khalil Selmi et Mme Ines Ben Nasr. ~30 s.",

    # 2 — Outline
    "Parcourir les 7 sections rapidement. Contexte, solutions existantes, "
    "notre proposition, méthodologie Scrum, conception, réalisation sprint "
    "par sprint, tests/démo, conclusion. ~30 s.",

    # 3 — Problem
    "Deux points d'ancrage : étudiants/travailleurs voyagent beaucoup mais "
    "ont peu de budget, et les voitures privées roulent à moitié vides. "
    "Lire la problématique en bas — c'est la thèse de tout le projet. ~45 s.",

    # 4 — Company
    "Court. Ils ont fourni l'encadrement technique ET le cadre agile. "
    "Ne pas s'attarder, le jury s'intéresse au projet. ~30 s.",

    # 5 — Existing solutions
    "Ne pas lire toute la table. Pointer une ligne : « Regardez la "
    "réduction étudiante — aucune ne l'a. Portefeuille — aucune ne l'a. "
    "C'est notre niche. » Finir avec la phrase en gras. ~45 s.",

    # 6 — Proposed solution
    "Les trois piliers structurent le reste de l'exposé. Chiffres à "
    "retenir : 3 applis, 4 sprints, 20+ entités, 3 langues. ~45 s.",

    # 7 — Methodology (THE slide for the supervisor)
    "C'EST la diapo qui intéresse l'encadrant. Parcourir la frise : Sprint 0 "
    "préparatoire, quatre sprints fonctionnels, groupés en deux releases. "
    "Mentionner les stand-ups quotidiens et les revues hebdomadaires. ~1 min.",

    # 8 — Architecture
    "Deux clients, un backend, une base. Les règles métier sont dans le "
    "backend, donc impossible à contourner depuis le mobile. Bearer tokens, "
    "bcrypt, HTTPS. ~1 min.",

    # 9 — Class diagram
    "NE PAS énumérer chaque classe. Dire : une vingtaine d'entités, "
    "organisées en cinq groupes logiques — utilisateurs, trajets, paiements, "
    "vérification, communication. Renvoyer au rapport pour le détail. ~30 s.",

    # 10 — Use case
    "Même principe : ne pas tout énumérer. « Quatre acteurs — passager, "
    "conducteur, admin, et le système lui-même pour les opérations "
    "automatiques comme le calcul de la réduction ou du score de fiabilité. » "
    "~30 s.",

    # 11 — Sprint 1
    "Lire l'objectif (rail gauche). Mettre l'accent sur US1.4 (OTP) et "
    "US1.5 (candidature conducteur) — les plus intéressants. Lire les "
    "3 livrables. ~1 min.",

    # 12 — Sprint 2
    "Objectif, puis pointer US2.6 (acompte 50 %) et US2.7 (réservation de "
    "groupe). Mentionner la collecte cash dans le diagramme d'activité — "
    "ça montre que vous avez pensé au flux complet, pas seulement au "
    "scénario nominal. ~1 min.",

    # 13 — Sprint 3
    "Insister : « Bug du pipeline de prix détecté en revue — la promo "
    "s'appliquait AVANT la réduction étudiante, légèrement en faveur du "
    "passager. Corrigé le soir même + test unitaire ajouté. » Ça prouve "
    "que vous faites de vraies revues. ~1 min.",

    # 14 — Sprint 4
    "Mentionner : le chat est en polling 3 secondes, sans WebSocket — un "
    "choix délibéré, pas une limite. Le bot HelpDesk couvre 14 catégories "
    "FAQ avec escalade. L'arabe bascule automatiquement la mise en page. "
    "~1 min. (Vous devriez être à ~12 min à la fin de cette diapo.)",

    # 15 — Selected screens
    "De gauche à droite : auth → recherche avec score → détail trajet avec "
    "carte → tableau de bord conducteur avec jauge de fiabilité. ~45 s.",

    # 16 — Mobile features
    "Si vous êtes en retard, SAUTER cette diapo — le jury peut lire. "
    "Sinon : push (FCM), carte interactive (OSM + OSRM), chat intégré, "
    "multilingue avec RTL automatique. ~30 s.",

    # 17 — Stack
    "Si en retard, SAUTER aussi. Une phrase par colonne : Web (PHP + "
    "Bootstrap), Backend (PHP REST + bearer + SQLite WAL), Mobile "
    "(Flutter + Provider + go_router + FCM). ~30 s.",

    # 18 — Tests
    "« Les tests étaient continus, pas reportés. » Mentionner la suite "
    "unitaire du pipeline de prix — c'est là que le bug du Sprint 3 a été "
    "verrouillé. ~45 s.",

    # 19 — LIVE DEMO
    "SCRIPT DÉMO (3 min, 11 étapes) : (1) Inscription passager sur mobile. "
    "(2) Taper un email universitaire reconnu → bannière verte. "
    "(3) Déclencher la vérification OTP. (4) Passer sur le web → session "
    "conducteur → publier un trajet Tunis-Sfax. (5) Mobile → rechercher "
    "Tunis Sfax → le trajet apparaît avec le badge de score. (6) Ouvrir le "
    "détail → carte + détail du prix avec la réduction étudiante. "
    "(7) Réserver une place → confirmation de l'acompte 50 %. (8) Web → "
    "le conducteur accepte. (9) Mobile → réservation confirmée dans Mes "
    "Trajets. (10) Ouvrir le chat → envoyer « On arrive dans 5 minutes ». "
    "(11) Tableau de bord conducteur → jauge de fiabilité. "
    "SECOURS : basculer sur le PDF des captures et narrer.",

    # 20 — Limitations
    "ÊTRE HONNÊTE — les jurys testent l'honnêteté ici. Reconnaître : "
    "portefeuille par rechargement manuel, SQLite pour le dev pas la "
    "prod, iOS non vérifié. Puis pivoter vers la feuille de route. ~1 min.",

    # 21 — Thanks
    "Nous vous remercions pour votre attention. Nous sommes prêts pour vos "
    "questions. ~15 s.",
]

for slide, note in zip(prs.slides, SPEAKER_NOTES):
    slide.notes_slide.notes_text_frame.text = note

out = os.path.join(ROOT, "ForsaDrive_Defense_FR.pptx")
prs.save(out)
print(f"Wrote {out}")
print(f"Slides: {len(prs.slides)}")
