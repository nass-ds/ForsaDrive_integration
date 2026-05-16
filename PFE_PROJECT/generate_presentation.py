#!/usr/bin/env python3
"""
ForsaDrive — defense presentation generator (~20 slides, 16:9).

Uses the same brand colors as the mobile app and references the same
diagrams that ship with the LaTeX report.

Run:
    python3 generate_presentation.py

Output:
    ForsaDrive_Defense.pptx (in the current directory)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────────────────────────────────────────────────────────────
# Brand
# ─────────────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x16, 0x28)   # primary, mobile/lib/utils/app_theme.dart
NAVY_DARK   = RGBColor(0x06, 0x0E, 0x1A)
ACCENT      = RGBColor(0xE8, 0xB8, 0x4B)   # gold
ACCENT_DARK = RGBColor(0xD4, 0xA4, 0x3A)
CREAM       = RGBColor(0xFF, 0xF3, 0xCD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF3, 0xF4, 0xF6)
DARK_TEXT   = RGBColor(0x1F, 0x2A, 0x3D)
MUTED       = RGBColor(0x6B, 0x72, 0x80)

ROOT = os.path.dirname(os.path.abspath(__file__))
def img(name): return os.path.join(ROOT, name)

# 16:9 dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=DARK_TEXT,
                bullet_color=ACCENT_DARK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet marker
        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.name = 'Calibri'
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        # body
        r = p.add_run()
        r.text = line
        r.font.name = 'Calibri'
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def page_bg(slide, color=WHITE):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def add_header_bar(slide, title, eyebrow=None):
    """Top header band with section eyebrow + slide title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), NAVY)
    add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.06), ACCENT)
    if eyebrow:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(8), Inches(0.35),
                 eyebrow.upper(), size=11, bold=True, color=ACCENT)
        add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                 title, size=26, bold=True, color=WHITE)
    else:
        add_text(slide, Inches(0.5), Inches(0.32), Inches(12), Inches(0.7),
                 title, size=28, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

def add_footer(slide, page_no, total, notes=None):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), NAVY)
    add_text(slide, Inches(0.5), Inches(7.2), Inches(8), Inches(0.3),
             "ForsaDrive  •  Final Year Project Defense  •  ATOMIC IT",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.5), Inches(7.2), Inches(1.5), Inches(0.3),
             f"{page_no} / {total}",
             size=10, color=ACCENT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
             bold=True)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_image_centered(slide, path, top, max_w, max_h):
    if not os.path.exists(path):
        # placeholder rectangle if image missing
        ph = add_rect(slide, (SLIDE_W - max_w) / 2, top, max_w, max_h, LIGHT_GREY,
                      line=MUTED)
        add_text(slide, (SLIDE_W - max_w) / 2, top + max_h / 2 - Inches(0.2),
                 max_w, Inches(0.4), f"[image: {os.path.basename(path)}]",
                 size=12, color=MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        return ph
    pic = slide.shapes.add_picture(path, 0, top, width=max_w, height=max_h)
    # Adjust to keep aspect ratio: shrink whichever dimension exceeds the box
    if pic.width / pic.height > max_w / max_h:
        # too wide: lock width
        pic.width = max_w
        pic.height = int(max_w * (pic.height / pic.width)) if False else pic.height
    # Re-center
    pic.left = int((SLIDE_W - pic.width) / 2)
    pic.top  = top
    return pic

def add_table(slide, x, y, w, h, headers, rows, *, header_fill=NAVY,
              header_color=WHITE, body_fill=WHITE, alt_fill=LIGHT_GREY,
              text_color=DARK_TEXT, font_size=12, header_size=12):
    cols  = len(headers)
    nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, cols, x, y, w, h)
    tbl = tbl_shape.table
    # Header
    for j, h_text in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h_text
        r.font.name = 'Calibri'
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_color
    # Body
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = body_fill if i % 2 == 0 else alt_fill
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = 'Calibri'
            r.font.size = Pt(font_size)
            r.font.color.rgb = text_color
    return tbl

def add_chip(slide, x, y, label, *, fill=ACCENT, text_color=NAVY, size=11):
    w = Inches(max(0.7, 0.12 * len(label) + 0.4))
    h = Inches(0.32)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = 'Calibri'
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shp, w

# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────
TOTAL = 21  # placeholder; recomputed at end

def slide_title():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    # decorative accent strip
    add_rect(s, 0, Inches(2.4), Inches(0.35), Inches(2.7), ACCENT)

    # Eyebrow
    add_text(s, Inches(1.0), Inches(2.3), Inches(11), Inches(0.4),
             "FINAL YEAR PROJECT  •  2025 — 2026", size=14, bold=True,
             color=ACCENT)
    # Big title
    add_text(s, Inches(1.0), Inches(2.7), Inches(11.5), Inches(1.2),
             "ForsaDrive", size=66, bold=True, color=WHITE)
    # Subtitle
    add_text(s, Inches(1.0), Inches(3.9), Inches(11.5), Inches(0.7),
             "A Tunisian-context carpooling platform — Web, Mobile, and a unified backend",
             size=20, color=CREAM)

    # Authors block
    add_rect(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.04), ACCENT)
    add_text(s, Inches(1.0), Inches(5.4), Inches(5.5), Inches(0.4),
             "Presented by", size=11, color=ACCENT, bold=True)
    add_text(s, Inches(1.0), Inches(5.65), Inches(5.5), Inches(0.5),
             "Youssef BEN ABID  &  Anas YOUNES", size=18, bold=True, color=WHITE)

    add_text(s, Inches(7.0), Inches(5.4), Inches(5.5), Inches(0.4),
             "Supervised by", size=11, color=ACCENT, bold=True)
    add_text(s, Inches(7.0), Inches(5.65), Inches(5.5), Inches(0.4),
             "Mr. Khalil SELMI  (ATOMIC IT)", size=14, color=WHITE)
    add_text(s, Inches(7.0), Inches(5.95), Inches(5.5), Inches(0.4),
             "Ms. Ines BEN NASR  (Academic)", size=14, color=WHITE)

    # Footer
    add_rect(s, 0, Inches(7.15), SLIDE_W, Inches(0.35), NAVY_DARK)
    add_text(s, Inches(0.5), Inches(7.2), Inches(12), Inches(0.3),
             "ATOMIC IT  •  Kelibia, Nabeul, Tunisia",
             size=10, color=ACCENT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

def slide_plan():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Outline of the defense", eyebrow="Plan")

    items = [
        ("01", "Project context & problem", "Why mobility in Tunisia is a real problem"),
        ("02", "Existing solutions", "What people use today and where it falls short"),
        ("03", "Proposed solution & methodology", "ForsaDrive concept and the Scrum process"),
        ("04", "Conception", "Architecture, use cases, class diagram"),
        ("05", "Realization — sprint by sprint", "From foundations to community features"),
        ("06", "Tests & demo", "What works and a live walk-through"),
        ("07", "Conclusion & future work", "What we delivered and what comes next"),
    ]
    y = Inches(1.7)
    for num, title, sub in items:
        add_rect(s, Inches(0.6), y, Inches(0.55), Inches(0.55), ACCENT)
        add_text(s, Inches(0.6), y, Inches(0.55), Inches(0.55),
                 num, size=16, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.4), y - Inches(0.02), Inches(11.5), Inches(0.35),
                 title, size=18, bold=True, color=NAVY)
        add_text(s, Inches(1.4), y + Inches(0.28), Inches(11.5), Inches(0.3),
                 sub, size=12, color=MUTED)
        y += Inches(0.72)

    add_footer(s, 2, TOTAL)

def slide_problem():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Mobility in Tunisia — a real, daily problem",
                   eyebrow="Context & motivation")

    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.5),
             "Public transport is irregular. Private cars travel half empty. "
             "Existing apps are not built for the local market.",
             size=15, color=DARK_TEXT)

    # Three columns of pain points
    cols = [
        ("Public transport",
         "Overcrowded buses, irregular schedules, weak coverage of secondary cities."),
        ("Informal carpooling",
         "Facebook groups, no identity verification, no payment trace, no rating."),
        ("International apps",
         "BlaBlaCar/inDrive/Bolt — payment models and language not adapted to Tunisia."),
    ]
    x = Inches(0.7)
    for title, body in cols:
        add_rect(s, x, Inches(2.4), Inches(4.0), Inches(2.6), LIGHT_GREY)
        add_rect(s, x, Inches(2.4), Inches(4.0), Inches(0.06), ACCENT)
        add_text(s, x + Inches(0.25), Inches(2.55), Inches(3.6), Inches(0.5),
                 title, size=16, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), Inches(3.05), Inches(3.6), Inches(1.8),
                 body, size=12, color=DARK_TEXT)
        x += Inches(4.15)

    # Quote band
    add_rect(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.4), NAVY)
    add_rect(s, Inches(0.7), Inches(5.3), Inches(0.1), Inches(1.4), ACCENT)
    add_text(s, Inches(1.0), Inches(5.45), Inches(11.5), Inches(0.4),
             "PROBLEM STATEMENT", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.9),
             "Tunisia lacks a secure, affordable digital ride-sharing solution "
             "adapted to local payment habits, available on web and mobile, "
             "and capable of building genuine trust between drivers and passengers.",
             size=14, color=WHITE)

    add_footer(s, 3, TOTAL)

def slide_company():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "ATOMIC IT — host company", eyebrow="Internship framework")

    add_bullets(s, Inches(0.7), Inches(1.7), Inches(7.5), Inches(4.5),
        [
            "IT engineering services: web and mobile development, software consulting.",
            "Modern stack: Linux/Windows, .NET, Java, Android/iOS, PHP.",
            "Located in Kelibia, Nabeul. Director: Mr. Khalil SELMI.",
            "Flat structure — short feedback loops between developers and management.",
            "Internship inside the engineering team, mentored by senior developers.",
            "Agile Scrum applied: daily stand-ups, weekly reviews with the supervisor.",
        ], size=15)

    # Right-side info card
    add_rect(s, Inches(8.7), Inches(1.7), Inches(4.0), Inches(4.5), LIGHT_GREY)
    add_rect(s, Inches(8.7), Inches(1.7), Inches(4.0), Inches(0.06), ACCENT)
    add_text(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(0.4),
             "AT A GLANCE", size=11, bold=True, color=NAVY)

    info = [
        ("Activity", "Web & mobile development"),
        ("Director", "Mr. Khalil SELMI"),
        ("Location", "Kelibia, Nabeul"),
        ("Phone", "+216 55 343 224"),
        ("Email", "contact@atomicitpro.com"),
        ("Methodology", "Agile / Scrum"),
    ]
    y = Inches(2.35)
    for k, v in info:
        add_text(s, Inches(8.95), y, Inches(1.3), Inches(0.32),
                 k, size=10, bold=True, color=MUTED)
        add_text(s, Inches(10.25), y, Inches(2.4), Inches(0.32),
                 v, size=11, color=DARK_TEXT)
        y += Inches(0.36)

    add_footer(s, 4, TOTAL)

def slide_existing():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "What people use today", eyebrow="Critical study")

    headers = ["Criterion", "Facebook", "inDrive", "Bolt", "BlaBlaCar", "ForsaDrive"]
    rows = [
        ["Identity verification",   "—",   "Partial", "✓", "✓", "✓"],
        ["In-app payment",          "—",   "—",       "✓", "✓", "✓"],
        ["Long-distance rides",     "✓",   "—",       "—", "✓", "✓"],
        ["Rating system",           "—",   "✓",       "✓", "✓", "✓"],
        ["Adapted to local context","✓",   "Partial", "Partial", "—", "✓"],
        ["Student discount (50%)",  "—",   "—",       "—", "—", "✓"],
        ["Wallet-based balance",    "—",   "—",       "—", "—", "✓"],
    ]
    add_table(s, Inches(0.7), Inches(1.7), Inches(12), Inches(4.6),
              headers, rows, font_size=12, header_size=12)

    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
             "No existing option combines local fit with structured trust mechanisms. "
             "ForsaDrive sits exactly in that gap.",
             size=13, bold=True, color=NAVY)

    add_footer(s, 5, TOTAL)

def slide_solution():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "ForsaDrive — what we built", eyebrow="Proposed solution")

    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.6),
        "An integrated ecosystem: web app, mobile app, and a shared backend, "
        "covering the full lifecycle from registration to post-trip rating.",
        size=14, color=DARK_TEXT)

    pillars = [
        ("Local fit",
         "50% prepayment + cash settlement, French/English/Arabic with RTL, "
         "Tunisian-recognized university domains."),
        ("Trust",
         "Identity & student verification, ratings on both sides, "
         "complaint workflow, reliability score."),
        ("Smart features",
         "Compatibility score, driver analytics, ride boosting, "
         "social feed, in-app chat, AI-assisted helpdesk."),
    ]
    x = Inches(0.7)
    for title, body in pillars:
        add_rect(s, x, Inches(2.5), Inches(4.0), Inches(3.0), NAVY)
        add_rect(s, x, Inches(2.5), Inches(4.0), Inches(0.08), ACCENT)
        add_text(s, x + Inches(0.25), Inches(2.65), Inches(3.6), Inches(0.5),
                 title, size=18, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.25), Inches(3.15), Inches(3.6), Inches(2.2),
                 body, size=13, color=WHITE)
        x += Inches(4.15)

    # KPI row
    kpis = [
        ("3", "coordinated apps"),
        ("4", "functional sprints"),
        ("20+", "domain entities"),
        ("3", "supported languages"),
    ]
    x = Inches(0.7)
    for v, lbl in kpis:
        add_rect(s, x, Inches(5.7), Inches(2.95), Inches(1.05), CREAM)
        add_text(s, x, Inches(5.75), Inches(2.95), Inches(0.55),
                 v, size=28, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(6.25), Inches(2.95), Inches(0.4),
                 lbl, size=11, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER)
        x += Inches(3.05)

    add_footer(s, 6, TOTAL)

def slide_methodology():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Scrum — five iterations grouped into two releases",
                   eyebrow="Methodology")

    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.5),
             "Sprint 0 prepared the ground. Four functional sprints then delivered "
             "two coherent releases of the product.",
             size=14, color=DARK_TEXT)

    # Sprint timeline (5 boxes)
    sprints = [
        ("S0", "Architecture\n& Conception", NAVY, ACCENT),
        ("S1", "Foundations\n(Auth + Profiles)", ACCENT, NAVY),
        ("S2", "Rides\n& Bookings", ACCENT, NAVY),
        ("S3", "Payments\n& Intelligence", ACCENT, NAVY),
        ("S4", "Community\n& Finalization", ACCENT, NAVY),
    ]
    box_w = Inches(2.3)
    gap   = Inches(0.15)
    total_w = box_w * 5 + gap * 4
    x = (SLIDE_W - total_w) / 2
    y = Inches(2.6)
    for code, title, fill, text_color in sprints:
        add_rect(s, x, y, box_w, Inches(2.0), fill)
        add_text(s, x, y + Inches(0.15), box_w, Inches(0.5),
                 code, size=22, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.85), box_w - Inches(0.2),
                 Inches(1.0), title, size=12, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += box_w + gap

    # Release labels
    add_rect(s, Inches(3.0), Inches(4.85), Inches(4.6), Inches(0.04), ACCENT)
    add_text(s, Inches(3.0), Inches(4.95), Inches(4.6), Inches(0.4),
             "RELEASE 1 — usable booking loop",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(7.7), Inches(4.85), Inches(4.6), Inches(0.04), ACCENT)
    add_text(s, Inches(7.7), Inches(4.95), Inches(4.6), Inches(0.4),
             "RELEASE 2 — intelligence + community",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Process bullets
    add_bullets(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.4),
        [
            "Daily stand-ups inside the team — weekly reviews with the supervising engineer.",
            "Each sprint produced a usable increment, reviewed and demonstrated.",
            "Trello board kept the backlog visible — story points on a Fibonacci scale.",
        ], size=13)

    add_footer(s, 7, TOTAL)

def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Architecture — client / server, three tiers",
                   eyebrow="Sprint 0 · Conception")

    # Left column — bullets
    add_bullets(s, Inches(0.7), Inches(1.7), Inches(5.4), Inches(5.0),
        [
            "Two clients (web + mobile) consume one REST API.",
            "Backend in PHP, exposed under /api/, secured with bearer tokens.",
            "Three tiers: presentation, application, data.",
            "External services (payment, notifications, file storage) only "
            "reached through the backend — no client secrets.",
            "HTTPS in production. Passwords with bcrypt. Audit log on sensitive "
            "operations.",
        ], size=14)

    # Right side — schematic boxes
    box_x = Inches(7.0); box_y = Inches(1.85)
    add_rect(s, box_x,             box_y,             Inches(2.6), Inches(0.9), CREAM)
    add_text(s, box_x, box_y, Inches(2.6), Inches(0.9), "Web App",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, box_x + Inches(3.0), box_y,             Inches(2.6), Inches(0.9), CREAM)
    add_text(s, box_x + Inches(3.0), box_y, Inches(2.6), Inches(0.9), "Mobile App",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, box_x + Inches(1.5), box_y + Inches(1.4), Inches(2.6), Inches(0.95), NAVY)
    add_text(s, box_x + Inches(1.5), box_y + Inches(1.4), Inches(2.6), Inches(0.95),
             "Backend\n(REST API)",
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, box_x + Inches(1.5), box_y + Inches(2.85), Inches(2.6), Inches(0.85), ACCENT)
    add_text(s, box_x + Inches(1.5), box_y + Inches(2.85), Inches(2.6), Inches(0.85),
             "SQLite (WAL)",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Connecting lines (visual hint)
    line1 = s.shapes.add_connector(2, box_x + Inches(1.3), box_y + Inches(0.9),
                                   box_x + Inches(2.5), box_y + Inches(1.4))
    line1.line.color.rgb = NAVY
    line2 = s.shapes.add_connector(2, box_x + Inches(4.3), box_y + Inches(0.9),
                                   box_x + Inches(3.0), box_y + Inches(1.4))
    line2.line.color.rgb = NAVY
    line3 = s.shapes.add_connector(2, box_x + Inches(2.8), box_y + Inches(2.35),
                                   box_x + Inches(2.8), box_y + Inches(2.85))
    line3.line.color.rgb = NAVY

    add_footer(s, 8, TOTAL)

def slide_class_diagram():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Class diagram — 20+ entities, 5 logical groups",
                   eyebrow="Sprint 0 · Conception")

    add_image_centered(s, img("forsadrive_class_diagram.png"),
                       Inches(1.4), Inches(11.5), Inches(4.95))

    # Caption / groups (above the footer band at 7.15)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             "Users & Profiles  ·  Trips & Bookings  ·  Payments & Ratings  "
             "·  Verification & Discounts  ·  Communication & Moderation",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(s, 9, TOTAL)

def slide_use_case():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Global use case — four actors", eyebrow="Sprint 0 · Conception")
    add_image_centered(s, img("ForsaDrive_UseCase.png"),
                       Inches(1.4), Inches(11.5), Inches(4.95))
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             "Passenger  ·  Driver  ·  Administrator  ·  System (automated)",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, 10, TOTAL)

def slide_sprint_section(num, color_dark, color_accent, title, mission,
                         backlog_rows, deliverables, page_no):
    """Section divider + content side by side, single slide per sprint."""
    s = prs.slides.add_slide(BLANK)
    page_bg(s)

    # Left rail (sprint marker)
    add_rect(s, 0, 0, Inches(3.4), SLIDE_H, color_dark)
    add_text(s, Inches(0.4), Inches(0.6), Inches(2.6), Inches(0.5),
             f"SPRINT {num}", size=14, bold=True, color=color_accent)
    add_text(s, Inches(0.4), Inches(1.0), Inches(2.6), Inches(2.0),
             title, size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.4), Inches(3.4), Inches(2.6), Inches(0.4),
             "MISSION", size=11, bold=True, color=color_accent)
    add_text(s, Inches(0.4), Inches(3.7), Inches(2.7), Inches(2.0),
             mission, size=12, color=WHITE)

    # Right area — backlog table
    add_text(s, Inches(3.7), Inches(0.4), Inches(9), Inches(0.5),
             "Sprint backlog (extract)", size=14, bold=True, color=NAVY)
    headers = ["#", "User story", "Pts"]
    add_table(s, Inches(3.7), Inches(0.85), Inches(9.3), Inches(3.3),
              headers, backlog_rows, font_size=11, header_size=11)

    # Deliverables row
    add_text(s, Inches(3.7), Inches(4.35), Inches(9), Inches(0.4),
             "Deliverables", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(3.7), Inches(4.75), Inches(9.3), Inches(2.3),
                deliverables, size=12)

    add_footer(s, page_no, TOTAL)

def slide_sprint1():
    slide_sprint_section(
        1, NAVY, ACCENT,
        "Foundations",
        "Account creation, authentication, profile, "
        "student verification, driver application & admin approval.",
        [
            ["US1.1", "Register an account",                        "5"],
            ["US1.2", "Login with email + password",                "3"],
            ["US1.3", "View / update profile",                      "3"],
            ["US1.4", "Student verification via university email OTP","8"],
            ["US1.5", "Apply to become a driver (license + vehicle)", "8"],
            ["US1.6", "Admin reviews driver applications",          "5"],
            ["US1.7", "Admin manages recognized university domains","3"],
        ],
        [
            "Self-service OTP flow (10-min code, 5-attempt cap, 30s cooldown).",
            "Driver application creates DriverProfile + Vehicle on approval.",
            "Audit log entry on every verification event.",
        ],
        11)

def slide_sprint2():
    slide_sprint_section(
        2, NAVY, ACCENT,
        "Rides & Bookings",
        "Driver publishes a ride. Passenger searches, filters, books with "
        "50% prepayment.",
        [
            ["US2.1", "Publish a ride (origin, destination, date, price, seats)", "5"],
            ["US2.2", "Modify or cancel a ride",                                  "3"],
            ["US2.3", "Search rides by origin / destination / date",              "5"],
            ["US2.4", "Filter by price, departure time, driver rating",           "5"],
            ["US2.5", "View ride detail and driver profile",                      "3"],
            ["US2.6", "Book seats with 50% prepayment",                           "8"],
            ["US2.7", "Group booking for friends",                                "5"],
            ["US2.8", "Driver accepts / rejects booking requests",                "5"],
        ],
        [
            "Booking persisted only if prepayment succeeds — no orphan bookings.",
            "Search response indexed on origin, destination, departure_date.",
            "Activity diagram covered cash collection at end of trip.",
        ],
        12)

def slide_sprint3():
    slide_sprint_section(
        3, NAVY, ACCENT,
        "Payments & Intelligence",
        "Wallet, promo codes, ride boosting, driver analytics, "
        "compatibility score.",
        [
            ["US3.1", "View wallet balance and history",                  "5"],
            ["US3.2", "Top up wallet",                                    "3"],
            ["US3.3", "Auto-apply 50% student discount on booking",       "3"],
            ["US3.4", "Apply organizational promo code",                  "3"],
            ["US3.5", "Boost a ride to top of search/feed",               "5"],
            ["US3.6", "Driver dashboard with reliability + revenue + ratings","8"],
            ["US3.7", "Compatibility score badge per ride",               "5"],
            ["US3.8", "Recompute reliability after every completed trip", "3"],
        ],
        [
            "Price pipeline: base × seats → student 50% → promo % → 50% prepayment.",
            "Reliability score: rating × completion × (1 − cancellation).",
            "Bug caught at sprint review: promo applied before student discount — fixed.",
        ],
        13)

def slide_sprint4():
    slide_sprint_section(
        4, NAVY, ACCENT,
        "Community & Finalization",
        "Chat, ratings, complaints, social feed, helpdesk, "
        "multilingual interface, admin panel.",
        [
            ["US4.1", "Real-time in-app chat (poll-based)",          "5"],
            ["US4.2", "Push notifications (FCM)",                    "3"],
            ["US4.3", "Social feed (posts, likes, comments)",        "8"],
            ["US4.4", "Rate the other party after a trip",           "3"],
            ["US4.5", "File a complaint after a trip",               "5"],
            ["US4.6", "AI-assisted HelpDesk with bot escalation",    "5"],
            ["US4.7", "Switch interface FR / EN / AR (RTL)",         "5"],
            ["US4.8", "Admin panel — 7 tabs",                        "8"],
        ],
        [
            "Chat polled every 3 s — sufficient for current scale, no WebSocket dep.",
            "HelpDesk bot covers 14 FAQ categories, escalates to human agent.",
            "Arabic flips the layout automatically through Flutter localizations.",
        ],
        14)

def slide_screens():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Selected screens", eyebrow="Realization")

    # 4-column grid of mobile screens
    captions = [
        ("Authentication", "mobile_auth.png"),
        ("Search + match score", "mobile_search.png"),
        ("Trip details + map", "mobile_trip_details.png"),
        ("Driver dashboard", "mobile_driver_dashboard.png"),
    ]
    x = Inches(0.5); y = Inches(1.6)
    cell_w = Inches(3.0)
    for caption, fname in captions:
        add_rect(s, x, y, cell_w, Inches(4.4), LIGHT_GREY)
        add_image_centered_into(s, img(fname), x + Inches(0.1), y + Inches(0.1),
                                cell_w - Inches(0.2), Inches(4.0))
        add_text(s, x, y + Inches(4.1), cell_w, Inches(0.3),
                 caption, size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        x += cell_w + Inches(0.13)

    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             "Web counterparts cover the same workflows for desktop users "
             "and the administration panel.",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(s, 15, TOTAL)

def add_image_centered_into(slide, path, x, y, w, h):
    """Like add_image_centered but inside an arbitrary cell box."""
    if not os.path.exists(path):
        ph = add_rect(slide, x, y, w, h, WHITE, line=MUTED)
        add_text(slide, x, y + h / 2 - Inches(0.2), w, Inches(0.4),
                 f"[{os.path.basename(path)}]", size=10, color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return ph
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    # Center inside cell
    pic.left = int(x + (w - pic.width) / 2)
    pic.top  = int(y + (h - pic.height) / 2)
    return pic

def slide_mobile_features():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "What only the mobile makes possible",
                   eyebrow="Mobile-specific features")

    cells = [
        ("Push notifications",
         "Firebase Cloud Messaging delivers booking, message, and ride alerts "
         "even when the app is in the background."),
        ("Interactive map",
         "OpenStreetMap tiles + OSRM routing draw the driving line between "
         "origin and destination on each ride detail screen."),
        ("Real-time chat",
         "Conversations linked to bookings — bubble UI, date dividers, read receipts. "
         "3-second polling, no WebSocket dependency."),
        ("Multilingual + RTL",
         "FR / EN / AR with automatic right-to-left layout flipping for Arabic, "
         "no manual override in any widget."),
    ]
    x = Inches(0.6); y = Inches(1.6)
    w = Inches(6.0); h = Inches(2.55)
    for i, (title, body) in enumerate(cells):
        col = i % 2; row = i // 2
        cx = Inches(0.6) + col * (Inches(6.15))
        cy = Inches(1.6) + row * (Inches(2.7))
        add_rect(s, cx, cy, w, h, LIGHT_GREY)
        add_rect(s, cx, cy, Inches(0.08), h, ACCENT)
        add_text(s, cx + Inches(0.25), cy + Inches(0.2), w - Inches(0.4),
                 Inches(0.5), title, size=18, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.25), cy + Inches(0.85), w - Inches(0.4),
                 Inches(1.6), body, size=12, color=DARK_TEXT)

    add_footer(s, 16, TOTAL)

def slide_stack():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Technology stack", eyebrow="Realization")

    # Three columns — frontend / backend / mobile
    blocks = [
        ("Web frontend", NAVY, [
            ("PHP 8 (server-rendered)", "Routing, sessions, page rendering."),
            ("HTML5 + Bootstrap 5", "Responsive layout, consistent components."),
            ("Vanilla JavaScript", "Form validation, AJAX, admin tab switching."),
        ]),
        ("Backend / data", ACCENT_DARK, [
            ("PHP REST API", "Front-controller under /api/, JSON responses."),
            ("PDO + bearer tokens", "Parameterized queries, 32-byte hex tokens."),
            ("SQLite (WAL mode)", "Shared file between web and mobile."),
        ]),
        ("Mobile", NAVY, [
            ("Flutter (Dart)", "Single codebase, native rendering on Android & iOS."),
            ("Provider + go_router", "State management + declarative navigation."),
            ("FCM, flutter_map, l10n", "Push, OSM map, FR/EN/AR with RTL."),
        ]),
    ]

    x = Inches(0.6); y = Inches(1.6)
    w = Inches(4.0); h = Inches(5.0)
    for title, head_color, items in blocks:
        add_rect(s, x, y, w, Inches(0.65), head_color)
        add_text(s, x, y, w, Inches(0.65), title,
                 size=16, bold=True,
                 color=WHITE if head_color != ACCENT_DARK else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, y + Inches(0.65), w, h - Inches(0.65), LIGHT_GREY)
        cy = y + Inches(0.85)
        for tech, sub in items:
            add_text(s, x + Inches(0.25), cy, w - Inches(0.4), Inches(0.35),
                     tech, size=13, bold=True, color=NAVY)
            add_text(s, x + Inches(0.25), cy + Inches(0.35),
                     w - Inches(0.4), Inches(0.7),
                     sub, size=11, color=DARK_TEXT)
            cy += Inches(1.25)
        x += Inches(4.2)

    add_footer(s, 17, TOTAL)

def slide_tests():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Testing — continuous, not deferred",
                   eyebrow="Quality")

    # Left: testing strategy
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.4),
             "Strategy", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.0),
        [
            "Unit tests on the price pipeline, prepayment, booking state machine.",
            "Postman collections replay end-to-end scenarios after each backend change.",
            "Manual cross-device tests on three Android phones.",
            "Weekly demo with the supervising engineer caught usability issues.",
        ], size=12)

    # Right: results table
    add_text(s, Inches(6.5), Inches(1.7), Inches(6.3), Inches(0.4),
             "Results", size=16, bold=True, color=NAVY)
    rows = [
        ["Registration + student OTP",        "✓ Passed"],
        ["Token-based session",                "✓ Passed"],
        ["Driver application + approval",      "✓ Passed"],
        ["Publish / cancel a ride",            "✓ Passed"],
        ["Book with 50% prepayment",           "✓ Passed"],
        ["Student discount + promo code",      "✓ Passed"],
        ["Wallet top-up + history",            "✓ Passed"],
        ["Ratings + complaints",               "✓ Passed"],
        ["Real-time chat + push (FCM)",        "✓ Passed"],
        ["Multilingual + RTL",                 "✓ Passed"],
    ]
    add_table(s, Inches(6.5), Inches(2.1), Inches(6.3), Inches(4.7),
              ["Scenario", "Result"], rows, font_size=11, header_size=11)

    add_footer(s, 18, TOTAL)

def slide_demo():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    add_rect(s, 0, Inches(3.4), Inches(0.4), Inches(0.8), ACCENT)
    add_text(s, Inches(1.0), Inches(2.5), Inches(11.5), Inches(0.5),
             "LIVE DEMO", size=18, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(3.0), Inches(11.5), Inches(1.5),
             "ForsaDrive in action", size=64, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.7), Inches(11.5), Inches(0.7),
             "Walk-through: registration → student verification → search → "
             "book → driver acceptance → cash collection.",
             size=18, color=CREAM)
    add_footer(s, 19, TOTAL)

def slide_limitations():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_header_bar(s, "Limitations & roadmap", eyebrow="Honest assessment")

    # Two columns
    add_rect(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.0), LIGHT_GREY)
    add_rect(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(0.06), ACCENT)
    add_text(s, Inches(0.85), Inches(1.75), Inches(5.5), Inches(0.5),
             "What's missing today", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(5.6), Inches(4.5),
        [
            "Wallet uses manual top-ups — no real payment gateway integrated yet.",
            "SQLite is suited for development, not high-traffic production.",
            "Mobile tested only on Android — iOS build not verified.",
            "No forgot-password email flow yet.",
            "Performance under concurrent load was not measured.",
        ], size=12, bullet_color=RGBColor(0xC0, 0x39, 0x2B))

    add_rect(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(5.0), CREAM)
    add_rect(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(0.06), NAVY)
    add_text(s, Inches(6.95), Inches(1.75), Inches(5.5), Inches(0.5),
             "Next steps", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(6.95), Inches(2.3), Inches(5.6), Inches(4.5),
        [
            "Migrate to PostgreSQL or MySQL for production.",
            "Integrate a Tunisian payment gateway (Konnect or Paymee).",
            "Deploy backend behind HTTPS on a VPS, automated backups.",
            "Publish the Flutter app on the Google Play Store.",
            "Add 2FA for driver and admin accounts + offline cache on mobile.",
        ], size=12)

    add_footer(s, 20, TOTAL)

def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY)
    add_rect(s, 0, Inches(3.0), Inches(0.4), Inches(1.5), ACCENT)
    add_text(s, Inches(1.0), Inches(2.6), Inches(11.5), Inches(0.5),
             "QUESTIONS & DISCUSSION", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(3.1), Inches(11.5), Inches(1.5),
             "Thank you for your attention.",
             size=56, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.7), Inches(11.5), Inches(0.6),
             "We are now ready for your questions.",
             size=20, color=CREAM)

    # Bottom names band
    add_rect(s, 0, Inches(6.5), SLIDE_W, Inches(0.65), NAVY_DARK)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.55),
             "Youssef BEN ABID  &  Anas YOUNES   ·   ForsaDrive   ·   2025 — 2026",
             size=14, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             bold=True)
    add_footer(s, 21, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────────────────────────────────────
slide_title()
slide_plan()
slide_problem()
slide_company()
slide_existing()
slide_solution()
slide_methodology()
slide_architecture()
slide_class_diagram()
slide_use_case()
slide_sprint1()
slide_sprint2()
slide_sprint3()
slide_sprint4()
slide_screens()
slide_mobile_features()
slide_stack()
slide_tests()
slide_demo()
slide_limitations()
slide_thanks()

# ─────────────────────────────────────────────────────────────────────────────
# Speaker notes (visible in PowerPoint presenter view, hidden from audience)
# ─────────────────────────────────────────────────────────────────────────────
SPEAKER_NOTES = [
    # 1 — Title
    "Bonjour, je suis Youssef, voici Anas. Nous allons vous présenter "
    "ForsaDrive, notre PFE réalisé chez ATOMIC IT sous la supervision de "
    "Mr. Khalil Selmi et Mme Ines Ben Nasr. ~30 s.",

    # 2 — Outline
    "Walk the 7 sections quickly. Contexte, solutions existantes, notre "
    "proposition, méthodologie Scrum, conception, réalisation sprint par "
    "sprint, tests/demo, conclusion. ~30 s.",

    # 3 — Problem
    "Two anchors: students/workers travel a lot but have small budgets, "
    "and private cars travel half empty. Read the problem statement quote "
    "at the bottom — it's the thesis of the whole project. ~45 s.",

    # 4 — Company
    "Brief. They provided technical supervision AND the agile process. "
    "Don't dwell, the jury cares about the project. ~30 s.",

    # 5 — Existing solutions
    "Don't read the whole table. Pick one row: 'Look at student discount — "
    "none of these have it. Wallet — none of these have it. That's our gap.' "
    "End with the bold sentence. ~45 s.",

    # 6 — Proposed solution
    "Three pillars structure the rest of the talk. Numbers to memorize: "
    "3 apps, 4 sprints, 20+ entities, 3 languages. ~45 s.",

    # 7 — Methodology (THE slide for the supervisor)
    "THIS is the slide your supervisor cares about. Walk the timeline: "
    "Sprint 0 was preparatory, four functional sprints, grouped into two "
    "releases. Mention daily stand-ups and weekly reviews. ~1 min.",

    # 8 — Architecture
    "Two clients, one backend, one DB. The business rules live in the "
    "backend so they cannot be bypassed from the mobile side. Bearer "
    "tokens, bcrypt, HTTPS. ~1 min.",

    # 9 — Class diagram
    "Do NOT read every class. Say: 20-something entities, organized into "
    "five logical groups — users, trips, payments, verification, "
    "communication. Refer the jury to the report for the detail. ~30 s.",

    # 10 — Use case
    "Same idea: don't enumerate. 'Four actors — passenger, driver, admin, "
    "and the system itself for automated operations like discount "
    "calculation and reliability scoring.' ~30 s.",

    # 11 — Sprint 1
    "Read the mission (left rail). Highlight US1.4 (OTP) and US1.5 (driver "
    "app) — they're the most interesting. Read the 3 deliverables. ~1 min.",

    # 12 — Sprint 2
    "Mission, then highlight US2.6 (50% prepayment) and US2.7 (group "
    "booking). Mention the cash collection step in the activity diagram — "
    "shows you thought about the full flow, not just the happy path. ~1 min.",

    # 13 — Sprint 3
    "Highlight: 'The price pipeline bug found at the sprint review — we "
    "applied the promo BEFORE the student discount, slightly favoring the "
    "passenger. Fixed same evening + unit test added.' Shows you do real "
    "reviews. ~1 min.",

    # 14 — Sprint 4
    "Mention: chat is 3-second polling, no WebSocket — a deliberate "
    "trade-off, not a limitation. HelpDesk bot covers 14 FAQ categories "
    "with escalation. Arabic flips the layout automatically. ~1 min. "
    "(You should be at ~12 min total by end of this slide.)",

    # 15 — Selected screens
    "Walk left to right: auth → search with match score → trip detail with "
    "map → driver dashboard with reliability gauge. ~45 s.",

    # 16 — Mobile features
    "If running long, SKIP this slide — the jury can read it. Otherwise: "
    "push (FCM), interactive map (OSM + OSRM), in-app chat, multilingual "
    "with auto RTL. ~30 s.",

    # 17 — Stack
    "If running long, SKIP this too. One sentence per column: Web (PHP + "
    "Bootstrap), Backend (PHP REST + bearer + SQLite WAL), Mobile (Flutter + "
    "Provider + go_router + FCM). ~30 s.",

    # 18 — Tests
    "'Tests were continuous, not deferred.' Mention the price-pipeline "
    "unit suite specifically — that's where the Sprint 3 bug was locked. "
    "~45 s.",

    # 19 — LIVE DEMO
    "DEMO SCRIPT (3 min, 11 steps): (1) Register passenger on mobile. "
    "(2) Type recognized university email → green banner. (3) Trigger OTP "
    "verification. (4) Switch to web → driver session → publish Tunis-Sfax "
    "ride. (5) Mobile → search Tunis Sfax → ride appears with match badge. "
    "(6) Open ride detail → map + price breakdown with student discount. "
    "(7) Book one seat → 50% prepayment confirmation. (8) Web → driver "
    "accepts. (9) Mobile → booking confirmed in My Rides. (10) Open chat → "
    "send 'On arrive dans 5 minutes'. (11) Show driver dashboard → "
    "reliability gauge. FALLBACK: switch to PDF of screenshots and narrate.",

    # 20 — Limitations
    "BE HONEST — jurors test honesty here. Acknowledge: wallet uses manual "
    "top-ups, SQLite is for dev not production, iOS not verified. Then "
    "pivot to the roadmap. ~1 min.",

    # 21 — Thanks
    "Nous vous remercions pour votre attention. Nous sommes prêts pour vos "
    "questions. ~15 s.",
]

for slide, note in zip(prs.slides, SPEAKER_NOTES):
    slide.notes_slide.notes_text_frame.text = note

out = os.path.join(ROOT, "ForsaDrive_Defense.pptx")
prs.save(out)
print(f"Wrote {out}")
print(f"Slides: {len(prs.slides)}")
